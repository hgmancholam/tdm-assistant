#!/usr/bin/env python3
"""
doc-import.py — Import documents into a project, convert to Markdown, categorize.

Supported formats: PDF, DOCX, XLSX, PPTX, CSV, TXT, MD
No AI is used for extraction. AI should only be called externally if the
document is a scanned PDF with no extractable text.

Usage:
    python doc-import.py import --file PATH --project CODE [--category auto] [--title TITLE]
    python doc-import.py list --project CODE
"""

import argparse
import csv
import json
import re
import sys
from datetime import datetime
from pathlib import Path

# ── Paths ─────────────────────────────────────────────────────────────────────

SCRIPT_DIR   = Path(__file__).resolve().parent          # .agents/skills/projects/
PROJECT_ROOT = SCRIPT_DIR.parents[2]                    # PersonalAssistant/

CATEGORIES = ["meetings", "decisions", "risks", "reports", "retrospectives"]

CATEGORY_KEYWORDS = {
    "meetings":       ["meeting", "minutes", "agenda", "standup", "sync", "kickoff",
                       "call", "weekly", "daily", "scrum", "1on1"],
    "decisions":      ["decision", "adr", "approved", "resolution", "record", "approval"],
    "risks":          ["risk", "risks", "issue", "blocker", "impediment", "mitigation"],
    "retrospectives": ["retro", "retrospective", "lessons", "learned", "postmortem", "post-mortem"],
    "reports":        ["report", "status", "summary", "update", "dashboard", "weekly-report"],
}

DEFAULT_CATEGORY = "reports"


# ── Category detection ─────────────────────────────────────────────────────────

def detect_category(filename: str) -> str:
    lower = filename.lower()
    for category, keywords in CATEGORY_KEYWORDS.items():
        if any(kw in lower for kw in keywords):
            return category
    return DEFAULT_CATEGORY


# ── Extractors ────────────────────────────────────────────────────────────────

def extract_pdf(path: Path) -> dict:
    """Extract text and tables from PDF. Uses pdfplumber (preferred) or pypdf (fallback)."""
    try:
        import pdfplumber
        pages = []
        with pdfplumber.open(str(path)) as pdf:
            for i, page in enumerate(pdf.pages, 1):
                text   = (page.extract_text() or "").strip()
                tables = page.extract_tables() or []
                pages.append({"page": i, "text": text, "tables": tables})
        extracted_text = sum(len(p["text"]) for p in pages)
        return {
            "format": "pdf",
            "title": _stem_to_title(path),
            "pages": pages,
            "extractable": extracted_text > 50,
        }
    except ImportError:
        pass

    try:
        import pypdf
        reader = pypdf.PdfReader(str(path))
        pages  = []
        for i, page in enumerate(reader.pages, 1):
            text = (page.extract_text() or "").strip()
            pages.append({"page": i, "text": text, "tables": []})
        extracted_text = sum(len(p["text"]) for p in pages)
        return {
            "format": "pdf",
            "title": _stem_to_title(path),
            "pages": pages,
            "extractable": extracted_text > 50,
        }
    except ImportError:
        raise ImportError("pdfplumber or pypdf is required to read PDF files. "
                          "Run: pip install pdfplumber")


def extract_docx(path: Path) -> dict:
    """Extract structured text and tables from Word documents (.docx only)."""
    if path.suffix.lower() == ".doc":
        raise ValueError(".doc (old Word format) is not supported. "
                         "Please save the file as .docx and try again.")
    try:
        from docx import Document
    except ImportError:
        raise ImportError("python-docx is required to read Word files. "
                          "Run: pip install python-docx")

    doc      = Document(str(path))
    title    = None
    sections = []
    cur_head = None
    cur_body = []

    for para in doc.paragraphs:
        style = para.style.name if para.style else ""
        text  = para.text.strip()
        if not text:
            continue
        if style.startswith("Heading 1") and not title:
            title = text
        if style.startswith("Heading"):
            if cur_head is not None or cur_body:
                sections.append({"heading": cur_head, "content": list(cur_body)})
            cur_head = text
            cur_body = []
        else:
            cur_body.append(text)

    if cur_head is not None or cur_body:
        sections.append({"heading": cur_head, "content": list(cur_body)})

    tables = []
    for tbl in doc.tables:
        rows = [[cell.text.strip() for cell in row.cells] for row in tbl.rows]
        if rows:
            tables.append(rows)

    return {
        "format": "docx",
        "title": title or _stem_to_title(path),
        "sections": sections,
        "tables": tables,
    }


def extract_xlsx(path: Path) -> dict:
    """Extract sheets and tables from Excel workbooks."""
    try:
        import openpyxl
    except ImportError:
        raise ImportError("openpyxl is required to read Excel files. "
                          "Run: pip install openpyxl")

    wb     = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    sheets = []
    for name in wb.sheetnames:
        ws   = wb[name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            if any(cell is not None for cell in row):
                rows.append([("" if cell is None else str(cell)) for cell in row])
        if rows:
            sheets.append({"name": name, "rows": rows})
    wb.close()

    return {
        "format": "xlsx",
        "title": _stem_to_title(path),
        "sheets": sheets,
    }


def extract_pptx(path: Path) -> dict:
    """Extract slides and text from PowerPoint presentations."""
    if path.suffix.lower() == ".ppt":
        raise ValueError(".ppt (old PowerPoint format) is not supported. "
                         "Please save the file as .pptx and try again.")
    try:
        from pptx import Presentation
        from pptx.util import Pt
    except ImportError:
        raise ImportError("python-pptx is required to read PowerPoint files. "
                          "Run: pip install python-pptx")

    prs    = Presentation(str(path))
    slides = []
    title  = None

    for i, slide in enumerate(prs.slides, 1):
        slide_title = None
        content     = []
        notes_text  = ""

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            # Detect title placeholder
            if hasattr(shape, "placeholder_format") and shape.placeholder_format:
                ph_type = shape.placeholder_format.type
                if ph_type == 1:   # TITLE
                    slide_title = text
                    if i == 1 and not title:
                        title = text
                    continue
            content.append(text)

        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()

        slides.append({
            "slide": i,
            "title": slide_title,
            "content": content,
            "notes": notes_text,
        })

    return {
        "format": "pptx",
        "title": title or _stem_to_title(path),
        "slides": slides,
    }


def extract_csv(path: Path) -> dict:
    """Extract rows from CSV file."""
    with open(str(path), encoding="utf-8-sig", errors="replace", newline="") as f:
        reader = csv.reader(f)
        rows   = [row for row in reader if any(cell.strip() for cell in row)]

    return {
        "format": "csv",
        "title": _stem_to_title(path),
        "sheets": [{"name": path.stem, "rows": rows}],
    }


def extract_text(path: Path) -> dict:
    """Extract plain text or Markdown files as-is."""
    content = path.read_text(encoding="utf-8", errors="replace")
    return {
        "format": "txt",
        "title": _stem_to_title(path),
        "text": content,
    }


# ── Markdown converters ────────────────────────────────────────────────────────

def _stem_to_title(path: Path) -> str:
    return re.sub(r"[-_]+", " ", path.stem).title()


def _rows_to_md_table(rows: list) -> str:
    if not rows:
        return ""
    header = rows[0]
    lines  = [
        "| " + " | ".join(str(c).replace("|", "\\|") for c in header) + " |",
        "| " + " | ".join("---" for _ in header) + " |",
    ]
    for row in rows[1:]:
        # Pad short rows to header length
        padded = list(row) + [""] * max(0, len(header) - len(row))
        lines.append("| " + " | ".join(str(c).replace("|", "\\|").replace("\n", " ")
                                       for c in padded) + " |")
    return "\n".join(lines)


def _header(data: dict, source: str, date: str, category: str) -> str:
    return (
        f"# {data['title']}\n\n"
        f"**Source:** `{source}`  \n"
        f"**Imported:** {date}  \n"
        f"**Category:** {category}\n\n"
    )


def to_markdown_pdf(data: dict, source: str, date: str, category: str) -> str:
    lines = [_header(data, source, date, category)]

    if not data.get("extractable", True):
        lines.append(
            "> **Warning:** This PDF appears to be scanned or image-based. "
            "No text could be extracted. Consider running OCR on the file first, "
            "or ask the TDM assistant to process it with AI."
        )
        return "\n".join(lines)

    for p in data.get("pages", []):
        if p["text"]:
            lines.append(p["text"])
            lines.append("")
        for table in p.get("tables", []):
            if table:
                lines.append(_rows_to_md_table(table))
                lines.append("")

    return "\n".join(lines)


def to_markdown_docx(data: dict, source: str, date: str, category: str) -> str:
    lines = [_header(data, source, date, category)]

    for section in data.get("sections", []):
        if section["heading"]:
            lines.append(f"## {section['heading']}\n")
        for para in section["content"]:
            lines.append(para)
            lines.append("")

    for i, table in enumerate(data.get("tables", []), 1):
        lines.append(f"### Table {i}\n")
        lines.append(_rows_to_md_table(table))
        lines.append("")

    return "\n".join(lines)


def to_markdown_tabular(data: dict, source: str, date: str, category: str) -> str:
    """For XLSX and CSV — renders each sheet as a Markdown table."""
    lines = [_header(data, source, date, category)]

    for sheet in data.get("sheets", []):
        lines.append(f"## {sheet['name']}\n")
        lines.append(_rows_to_md_table(sheet["rows"]))
        lines.append("")

    return "\n".join(lines)


def to_markdown_pptx(data: dict, source: str, date: str, category: str) -> str:
    lines = [_header(data, source, date, category)]

    for slide in data.get("slides", []):
        slide_title = slide.get("title") or f"Slide {slide['slide']}"
        lines.append(f"## Slide {slide['slide']}: {slide_title}\n")
        for content in slide.get("content", []):
            lines.append(content)
            lines.append("")
        if slide.get("notes"):
            lines.append(f"*Notes: {slide['notes']}*")
            lines.append("")

    return "\n".join(lines)


def to_markdown_text(data: dict, source: str, date: str, category: str) -> str:
    lines = [_header(data, source, date, category)]
    lines.append(data.get("text", ""))
    return "\n".join(lines)


# ── Dispatch tables ────────────────────────────────────────────────────────────

EXTRACTORS = {
    ".pdf":  extract_pdf,
    ".docx": extract_docx,
    ".doc":  extract_docx,
    ".xlsx": extract_xlsx,
    ".xls":  extract_xlsx,
    ".pptx": extract_pptx,
    ".ppt":  extract_pptx,
    ".csv":  extract_csv,
    ".txt":  extract_text,
    ".md":   extract_text,
}

CONVERTERS = {
    "pdf":  to_markdown_pdf,
    "docx": to_markdown_docx,
    "xlsx": to_markdown_tabular,
    "csv":  to_markdown_tabular,
    "pptx": to_markdown_pptx,
    "txt":  to_markdown_text,
}


# ── Commands ──────────────────────────────────────────────────────────────────

def cmd_import(file_path: str, project_code: str, category: str, title_override: str | None):
    path = Path(file_path).resolve()

    if not path.exists():
        return {"success": False, "error": f"File not found: {file_path}"}

    ext = path.suffix.lower()
    if ext not in EXTRACTORS:
        supported = ", ".join(sorted(EXTRACTORS))
        return {"success": False, "error": f"Unsupported format: {ext}. Supported: {supported}"}

    project_dir = PROJECT_ROOT / "projects" / project_code
    if not project_dir.exists():
        return {"success": False, "error": f"Project '{project_code}' not found in projects/. "
                                           f"Create it first with /new-project."}

    if category == "auto":
        category = detect_category(path.name)

    # Extract content (no AI)
    try:
        data = EXTRACTORS[ext](path)
    except (ImportError, ValueError) as e:
        return {"success": False, "error": str(e)}
    except Exception as e:
        return {"success": False, "error": f"Extraction failed: {type(e).__name__}: {e}"}

    if title_override:
        data["title"] = title_override

    # Convert to Markdown
    today     = datetime.now().strftime("%Y-%m-%d")
    fmt       = data["format"]
    converter = CONVERTERS.get(fmt, to_markdown_text)
    markdown  = converter(data, path.name, today, category)

    # Write to project folder
    out_dir = project_dir / category
    out_dir.mkdir(parents=True, exist_ok=True)

    slug     = re.sub(r"[^\w\-]", "-", path.stem.lower()).strip("-")
    out_path = out_dir / f"{today}-{slug}.md"

    counter = 1
    while out_path.exists():
        out_path = out_dir / f"{today}-{slug}-{counter}.md"
        counter += 1

    out_path.write_text(markdown, encoding="utf-8")

    return {
        "success":  True,
        "project":  project_code,
        "source":   str(path),
        "output":   str(out_path.relative_to(PROJECT_ROOT)),
        "category": category,
        "title":    data["title"],
        "format":   fmt,
        "scanned_pdf": fmt == "pdf" and not data.get("extractable", True),
    }


def cmd_list(project_code: str):
    project_dir = PROJECT_ROOT / "projects" / project_code
    if not project_dir.exists():
        return {"success": False, "error": f"Project '{project_code}' not found."}

    docs = []
    for cat in CATEGORIES:
        cat_dir = project_dir / cat
        if cat_dir.exists():
            for f in sorted(cat_dir.glob("*.md")):
                docs.append({
                    "category": cat,
                    "file":     f.name,
                    "path":     str(f.relative_to(PROJECT_ROOT)),
                    "size_kb":  round(f.stat().st_size / 1024, 1),
                })

    return {"success": True, "project": project_code, "count": len(docs), "documents": docs}


# ── CLI ───────────────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(
        description="Import documents into a project and convert to Markdown"
    )
    sub = parser.add_subparsers(dest="command", required=True)

    p_import = sub.add_parser("import", help="Import a document")
    p_import.add_argument("--file",     required=True, help="Path to the source document")
    p_import.add_argument("--project",  required=True, help="Project code (e.g. ALPHA)")
    p_import.add_argument("--category", default="auto",
                          choices=["auto"] + CATEGORIES,
                          help="Target folder (default: auto-detect from filename)")
    p_import.add_argument("--title", default=None, help="Override document title in Markdown")

    p_list = sub.add_parser("list", help="List imported documents for a project")
    p_list.add_argument("--project", required=True, help="Project code")

    args = parser.parse_args()

    if args.command == "import":
        result = cmd_import(args.file, args.project, args.category, args.title)
    else:
        result = cmd_list(args.project)

    print(json.dumps(result, indent=2, ensure_ascii=False))


if __name__ == "__main__":
    main()
